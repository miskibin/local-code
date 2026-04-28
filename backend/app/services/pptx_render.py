from __future__ import annotations

import base64
import tempfile
from collections.abc import Callable
from pathlib import Path
from typing import Any

from PIL import Image as PILImage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from app.schemas.pptx import (
    BulletsSlide,
    ChartSlide,
    ConclusionSlide,
    DeckSpec,
    SectionSlide,
    TableSlide,
    TitleSlide,
)

NAVY_HEX = "0F172A"
NAVY_SOFT_HEX = "1E293B"
ACCENT_HEX = "3B82F6"
ACCENT_DEEP_HEX = "1D4ED8"
INK_HEX = "0F172A"
INK_2_HEX = "334155"
MUTED_HEX = "64748B"
HAIRLINE_HEX = "E2E8F0"
SOFT_BG_HEX = "F8FAFC"
ROW_ALT_HEX = "F1F5F9"
WHITE_HEX = "FFFFFF"
WHITE_70_HEX = "CBD5E1"

PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)
MARGIN_X = Inches(0.65)
MARGIN_TOP = Inches(0.55)
MARGIN_BOTTOM = Inches(0.55)
TOP_STRIP_HEIGHT = Inches(0.08)
FOOTER_TOP = PAGE_H - Inches(0.45)


class PptxRenderError(RuntimeError):
    pass


def _layout_by_name(prs: Presentation, name: str):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    available = [layout.name for layout in prs.slide_layouts]
    raise PptxRenderError(f"template missing layout '{name}'; available: {available}")


def _add_text(
    slide,
    text: str,
    *,
    left,
    top,
    width,
    height,
    font_size: int = 18,
    bold: bool = False,
    color_hex: str = INK_HEX,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    letter_spacing: int | None = None,
    font_name: str = "Calibri",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = RGBColor.from_string(color_hex)
    return box


def _add_filled_rect(
    slide, *, left, top, width, height, fill_hex: str, line_hex: str | None = None
):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if line_hex is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = RGBColor.from_string(line_hex)
        rect.line.width = Pt(0.5)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    rect.shadow.inherit = False
    return rect


def _add_full_bleed_bg(slide, color_hex: str) -> None:
    _add_filled_rect(slide, left=0, top=0, width=PAGE_W, height=PAGE_H, fill_hex=color_hex)


def _add_top_strip(slide, color_hex: str = NAVY_HEX) -> None:
    _add_filled_rect(
        slide, left=0, top=0, width=PAGE_W, height=TOP_STRIP_HEIGHT, fill_hex=color_hex
    )


_TITLE_UNDERLINE_W = Inches(0.85)


def _add_title_underline(slide, *, top, color_hex: str = ACCENT_HEX, width=None) -> None:
    _add_filled_rect(
        slide,
        left=MARGIN_X,
        top=top,
        width=width if width is not None else _TITLE_UNDERLINE_W,
        height=Pt(3),
        fill_hex=color_hex,
    )


def _add_footer(slide, *, slide_idx: int, total_slides: int, deck_title: str | None) -> None:
    """Hairline + page number bottom-right, optional deck title bottom-left."""
    _add_filled_rect(
        slide,
        left=MARGIN_X,
        top=FOOTER_TOP - Inches(0.1),
        width=PAGE_W - MARGIN_X * 2,
        height=Pt(0.5),
        fill_hex=HAIRLINE_HEX,
    )
    if deck_title:
        _add_text(
            slide,
            deck_title.upper(),
            left=MARGIN_X,
            top=FOOTER_TOP,
            width=Inches(8),
            height=Inches(0.4),
            font_size=9,
            bold=False,
            color_hex=MUTED_HEX,
            letter_spacing=1,
        )
    _add_text(
        slide,
        f"{slide_idx:02d} / {total_slides:02d}",
        left=PAGE_W - MARGIN_X - Inches(2),
        top=FOOTER_TOP,
        width=Inches(2),
        height=Inches(0.4),
        font_size=10,
        bold=True,
        color_hex=MUTED_HEX,
        align=PP_ALIGN.RIGHT,
    )


def _add_bullets(
    slide,
    bullets: list[str],
    *,
    left,
    top,
    width,
    height,
    font_size: int = 20,
    marker_hex: str = ACCENT_HEX,
    text_hex: str = INK_2_HEX,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
        p.text = ""
        marker = p.add_run()
        marker.text = "▪  "
        marker.font.size = Pt(font_size)
        marker.font.bold = True
        marker.font.name = "Calibri"
        marker.font.color.rgb = RGBColor.from_string(marker_hex)
        body = p.add_run()
        body.text = b
        body.font.size = Pt(font_size)
        body.font.bold = False
        body.font.name = "Calibri"
        body.font.color.rgb = RGBColor.from_string(text_hex)
    return box


def _render_title(
    spec: TitleSlide, prs: Presentation, *, slide_idx, total_slides, deck_title
) -> None:
    layout = _layout_by_name(prs, "title")
    slide = prs.slides.add_slide(layout)
    _add_full_bleed_bg(slide, NAVY_HEX)
    _add_filled_rect(slide, left=0, top=0, width=Inches(0.18), height=PAGE_H, fill_hex=ACCENT_HEX)
    body_w = PAGE_W - MARGIN_X * 2 - Inches(0.18)
    left_x = MARGIN_X + Inches(0.18)
    _add_text(
        slide,
        "PRESENTATION",
        left=left_x,
        top=Inches(2.2),
        width=body_w,
        height=Inches(0.4),
        font_size=12,
        bold=True,
        color_hex=ACCENT_HEX,
        letter_spacing=2,
    )
    _add_filled_rect(
        slide,
        left=left_x,
        top=Inches(2.65),
        width=Inches(1.0),
        height=Pt(3),
        fill_hex=ACCENT_HEX,
    )
    _add_text(
        slide,
        spec.title,
        left=left_x,
        top=Inches(2.9),
        width=body_w,
        height=Inches(2.2),
        font_size=54,
        bold=True,
        color_hex=WHITE_HEX,
    )
    if spec.subtitle:
        _add_text(
            slide,
            spec.subtitle,
            left=left_x,
            top=Inches(5.1),
            width=body_w,
            height=Inches(1.2),
            font_size=22,
            bold=False,
            color_hex=WHITE_70_HEX,
        )


def _render_section(
    spec: SectionSlide, prs: Presentation, *, slide_idx, total_slides, deck_title
) -> None:
    layout = _layout_by_name(prs, "section")
    slide = prs.slides.add_slide(layout)
    _add_full_bleed_bg(slide, NAVY_HEX)
    _add_text(
        slide,
        f"{slide_idx:02d}",
        left=MARGIN_X,
        top=Inches(0.4),
        width=Inches(6),
        height=Inches(4),
        font_size=240,
        bold=True,
        color_hex=NAVY_SOFT_HEX,
    )
    body_w = PAGE_W - MARGIN_X * 2
    if spec.eyebrow:
        _add_text(
            slide,
            spec.eyebrow.upper(),
            left=MARGIN_X,
            top=Inches(3.3),
            width=body_w,
            height=Inches(0.5),
            font_size=13,
            bold=True,
            color_hex=ACCENT_HEX,
            letter_spacing=2,
        )
    _add_filled_rect(
        slide,
        left=MARGIN_X,
        top=Inches(3.85),
        width=Inches(0.85),
        height=Pt(3),
        fill_hex=ACCENT_HEX,
    )
    _add_text(
        slide,
        spec.title,
        left=MARGIN_X,
        top=Inches(4.1),
        width=body_w,
        height=Inches(2.0),
        font_size=48,
        bold=True,
        color_hex=WHITE_HEX,
    )


def _content_chrome(
    slide, *, slide_idx: int, total_slides: int, deck_title: str | None, title: str
) -> None:
    _add_top_strip(slide, NAVY_HEX)
    _add_text(
        slide,
        title,
        left=MARGIN_X,
        top=MARGIN_TOP,
        width=PAGE_W - MARGIN_X * 2,
        height=Inches(0.85),
        font_size=28,
        bold=True,
        color_hex=NAVY_HEX,
    )
    _add_title_underline(slide, top=MARGIN_TOP + Inches(0.85))
    _add_footer(slide, slide_idx=slide_idx, total_slides=total_slides, deck_title=deck_title)


def _render_bullets(
    spec: BulletsSlide, prs: Presentation, *, slide_idx, total_slides, deck_title
) -> None:
    layout = _layout_by_name(prs, "bullets")
    slide = prs.slides.add_slide(layout)
    _content_chrome(
        slide,
        slide_idx=slide_idx,
        total_slides=total_slides,
        deck_title=deck_title,
        title=spec.title,
    )
    _add_bullets(
        slide,
        spec.bullets,
        left=MARGIN_X,
        top=MARGIN_TOP + Inches(1.4),
        width=PAGE_W - MARGIN_X * 2,
        height=Inches(4.8),
        font_size=22,
    )


def _fit_picture_into_box(slide, image_path: str, *, left, top, width, height) -> None:
    with PILImage.open(image_path) as im:
        iw, ih = im.size
    if iw <= 0 or ih <= 0:
        raise PptxRenderError("image has zero dimension")
    img_aspect = iw / ih
    box_aspect = width / height
    if img_aspect > box_aspect:
        new_w = width
        new_h = int(width / img_aspect)
    else:
        new_h = height
        new_w = int(height * img_aspect)
    x = left + (width - new_w) // 2
    y = top + (height - new_h) // 2
    slide.shapes.add_picture(image_path, Emu(x), Emu(y), width=Emu(new_w), height=Emu(new_h))


def _render_chart(
    spec: ChartSlide, prs: Presentation, artifact_lookup, *, slide_idx, total_slides, deck_title
) -> None:
    artifact = artifact_lookup(spec.artifact_id)
    if artifact is None:
        raise PptxRenderError(f"artifact {spec.artifact_id!r} not found")
    if artifact.kind != "image":
        raise PptxRenderError(
            f"chart slide expects image artifact; {spec.artifact_id} is kind={artifact.kind!r}"
        )
    payload = artifact.payload or {}
    b64 = payload.get("data_b64") or payload.get("_image_png_b64")
    if not b64:
        raise PptxRenderError(f"image artifact {spec.artifact_id} payload missing data_b64")

    layout = _layout_by_name(prs, "chart")
    slide = prs.slides.add_slide(layout)
    _add_full_bleed_bg(slide, SOFT_BG_HEX)
    _content_chrome(
        slide,
        slide_idx=slide_idx,
        total_slides=total_slides,
        deck_title=deck_title,
        title=spec.title,
    )
    body_w = PAGE_W - MARGIN_X * 2
    img_top = MARGIN_TOP + Inches(1.4)
    cap_height = Inches(0.6) if spec.caption else Inches(0)
    img_height = FOOTER_TOP - Inches(0.2) - img_top - cap_height
    _add_filled_rect(
        slide,
        left=MARGIN_X,
        top=img_top,
        width=body_w,
        height=img_height + cap_height,
        fill_hex=WHITE_HEX,
        line_hex=HAIRLINE_HEX,
    )
    img_bytes = base64.b64decode(b64)
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        tmp.write(img_bytes)
        tmp_path = tmp.name
    try:
        _fit_picture_into_box(
            slide,
            tmp_path,
            left=MARGIN_X + Inches(0.2),
            top=img_top + Inches(0.2),
            width=body_w - Inches(0.4),
            height=img_height - Inches(0.2),
        )
    finally:
        Path(tmp_path).unlink(missing_ok=True)
    if spec.caption:
        _add_text(
            slide,
            spec.caption,
            left=MARGIN_X + Inches(0.3),
            top=img_top + img_height,
            width=body_w - Inches(0.6),
            height=cap_height,
            font_size=12,
            bold=False,
            color_hex=MUTED_HEX,
            align=PP_ALIGN.CENTER,
        )


def _parse_columns(columns_raw) -> tuple[list[str], list[str]]:
    keys: list[str] = []
    labels: list[str] = []
    for c in columns_raw:
        if isinstance(c, dict):
            keys.append(str(c.get("key") or c.get("label") or ""))
            labels.append(str(c.get("label") or c.get("key") or ""))
        else:
            keys.append(str(c))
            labels.append(str(c))
    return keys, labels


def _style_header_cell(cell, label: str) -> None:
    cell.text = label
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(11)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor.from_string(WHITE_HEX)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(NAVY_HEX)


def _style_body_cell(cell, value, *, zebra: bool) -> None:
    cell.text = "" if value is None else str(value)
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(10)
        run.font.name = "Calibri"
        run.font.color.rgb = RGBColor.from_string(INK_2_HEX)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(ROW_ALT_HEX if zebra else WHITE_HEX)


def _render_table(
    spec: TableSlide, prs: Presentation, artifact_lookup, *, slide_idx, total_slides, deck_title
) -> None:
    artifact = artifact_lookup(spec.artifact_id)
    if artifact is None:
        raise PptxRenderError(f"artifact {spec.artifact_id!r} not found")
    if artifact.kind != "table":
        raise PptxRenderError(
            f"table slide expects table artifact; {spec.artifact_id} is kind={artifact.kind!r}"
        )
    payload = artifact.payload or {}
    column_keys, column_labels = _parse_columns(payload.get("columns") or [])
    if not column_keys:
        raise PptxRenderError(f"table artifact {spec.artifact_id} has no columns")
    rows_raw = payload.get("rows") or []
    total_rows = len(rows_raw)
    visible = rows_raw[: spec.max_rows]
    truncated = total_rows > spec.max_rows

    layout = _layout_by_name(prs, "table")
    slide = prs.slides.add_slide(layout)
    _content_chrome(
        slide,
        slide_idx=slide_idx,
        total_slides=total_slides,
        deck_title=deck_title,
        title=spec.title,
    )
    body_w = PAGE_W - MARGIN_X * 2
    table_top = MARGIN_TOP + Inches(1.4)
    cap_height = Inches(0.45)
    table_height = FOOTER_TOP - Inches(0.2) - table_top - cap_height
    table_shape = slide.shapes.add_table(
        len(visible) + 1, len(column_keys), MARGIN_X, table_top, body_w, table_height
    )
    table = table_shape.table
    for j, label in enumerate(column_labels):
        _style_header_cell(table.cell(0, j), label)
    for i, row in enumerate(visible, start=1):
        for j, key in enumerate(column_keys):
            value = row.get(key) if isinstance(row, dict) else None
            _style_body_cell(table.cell(i, j), value, zebra=(i % 2 == 0))

    parts: list[str] = []
    if truncated:
        parts.append(f"Showing {len(visible)} of {total_rows} rows")
    if spec.caption:
        parts.append(spec.caption)
    if parts:
        _add_text(
            slide,
            " · ".join(parts),
            left=MARGIN_X,
            top=table_top + table_height + Inches(0.05),
            width=body_w,
            height=cap_height,
            font_size=11,
            bold=False,
            color_hex=MUTED_HEX,
        )


def _render_conclusion(
    spec: ConclusionSlide,
    prs: Presentation,
    *,
    slide_idx,
    total_slides,
    deck_title,
) -> None:
    layout = _layout_by_name(prs, "conclusion")
    slide = prs.slides.add_slide(layout)
    _content_chrome(
        slide,
        slide_idx=slide_idx,
        total_slides=total_slides,
        deck_title=deck_title,
        title=spec.title,
    )
    body_w = PAGE_W - MARGIN_X * 2
    bullets_height = Inches(3.6) if spec.cta else Inches(4.6)
    _add_bullets(
        slide,
        spec.bullets,
        left=MARGIN_X,
        top=MARGIN_TOP + Inches(1.4),
        width=body_w,
        height=bullets_height,
        font_size=22,
    )
    if spec.cta:
        cta_top = MARGIN_TOP + Inches(1.4) + bullets_height + Inches(0.3)
        cta_height = Inches(0.7)
        _add_filled_rect(
            slide,
            left=MARGIN_X,
            top=cta_top,
            width=body_w,
            height=cta_height,
            fill_hex=NAVY_HEX,
        )
        _add_filled_rect(
            slide,
            left=MARGIN_X,
            top=cta_top,
            width=Inches(0.18),
            height=cta_height,
            fill_hex=ACCENT_HEX,
        )
        _add_text(
            slide,
            spec.cta,
            left=MARGIN_X + Inches(0.4),
            top=cta_top,
            width=body_w - Inches(0.4),
            height=cta_height,
            font_size=15,
            bold=True,
            color_hex=WHITE_HEX,
        )


def render_deck(
    spec: DeckSpec,
    template_path: Path,
    artifact_lookup: Callable[[str], Any],
    out_path: Path,
) -> Path:
    template_path = Path(template_path)
    if not template_path.is_file():
        raise PptxRenderError(f"template not found: {template_path}")
    prs = Presentation(str(template_path))
    total = len(spec.slides)
    deck_title = spec.title
    for i, slide_spec in enumerate(spec.slides, start=1):
        kwargs = dict(slide_idx=i, total_slides=total, deck_title=deck_title)
        if isinstance(slide_spec, TitleSlide):
            _render_title(slide_spec, prs, **kwargs)
        elif isinstance(slide_spec, SectionSlide):
            _render_section(slide_spec, prs, **kwargs)
        elif isinstance(slide_spec, BulletsSlide):
            _render_bullets(slide_spec, prs, **kwargs)
        elif isinstance(slide_spec, ChartSlide):
            _render_chart(slide_spec, prs, artifact_lookup, **kwargs)
        elif isinstance(slide_spec, TableSlide):
            _render_table(slide_spec, prs, artifact_lookup, **kwargs)
        elif isinstance(slide_spec, ConclusionSlide):
            _render_conclusion(slide_spec, prs, **kwargs)
        else:
            raise PptxRenderError(f"no renderer for {type(slide_spec).__name__}")
    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


__all__ = [
    "PptxRenderError",
    "render_deck",
]
