from __future__ import annotations

import base64
from pathlib import Path
from types import SimpleNamespace

import pytest
from pptx import Presentation
from pydantic import ValidationError

from app.schemas.pptx import (
    BulletsSlide,
    ChartSlide,
    ConclusionSlide,
    DeckSpec,
    SectionSlide,
    TableSlide,
    TitleSlide,
)
from app.services.pptx_render import PptxRenderError, render_deck

TEMPLATE_PATH = (
    Path(__file__).resolve().parents[1] / "data" / "pptx_templates" / "default.pptx"
)


def _png_bytes() -> bytes:
    """1x1 transparent PNG."""
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNkAAIAAAoAAv/"
        "lxKUAAAAASUVORK5CYII="
    )


def _make_image_artifact(aid: str = "art_image1"):
    b64 = base64.b64encode(_png_bytes()).decode()
    return SimpleNamespace(
        id=aid,
        kind="image",
        title="img",
        payload={"format": "png", "data_b64": b64},
    )


def _make_table_artifact(aid: str = "art_table1", n_rows: int = 3):
    columns = [
        {"key": "name", "label": "Name"},
        {"key": "value", "label": "Value"},
    ]
    rows = [{"name": f"row{i}", "value": i * 10} for i in range(n_rows)]
    return SimpleNamespace(
        id=aid,
        kind="table",
        title="tbl",
        payload={"columns": columns, "rows": rows},
    )


def _make_text_artifact(aid: str = "art_text1"):
    return SimpleNamespace(
        id=aid, kind="text", title="t", payload={"text": "hello"}
    )


def _lookup(*artifacts):
    by_id = {a.id: a for a in artifacts}

    def _f(aid: str):
        return by_id.get(aid)

    return _f


def test_render_title_only_deck(tmp_path):
    spec = DeckSpec(
        title="Test Deck",
        slides=[
            TitleSlide(type="title", title="Hello world", subtitle="A subtitle"),
        ],
    )
    out = tmp_path / "out.pptx"
    path = render_deck(spec, TEMPLATE_PATH, _lookup(), out)
    assert path.is_file()
    prs = Presentation(str(path))
    assert len(prs.slides) == 1
    text_runs = [
        r.text
        for s in prs.slides
        for shape in s.shapes
        if shape.has_text_frame
        for p in shape.text_frame.paragraphs
        for r in p.runs
    ]
    assert any("Hello world" in t for t in text_runs)
    assert any("A subtitle" in t for t in text_runs)


def test_render_with_image_artifact(tmp_path):
    art = _make_image_artifact()
    spec = DeckSpec(
        title="Deck",
        slides=[
            ChartSlide(
                type="chart",
                title="Q3 revenue",
                artifact_id=art.id,
                caption="Up 18% YoY.",
            ),
        ],
    )
    out = tmp_path / "out.pptx"
    render_deck(spec, TEMPLATE_PATH, _lookup(art), out)
    prs = Presentation(str(out))
    slide = prs.slides[0]
    assert any(getattr(s, "image", None) is not None for s in slide.shapes), (
        "expected at least one Picture shape"
    )


def test_render_with_table_artifact(tmp_path):
    art = _make_table_artifact(n_rows=3)
    spec = DeckSpec(
        title="Deck",
        slides=[
            TableSlide(
                type="table",
                title="Top customers",
                artifact_id=art.id,
                caption="snapshot",
            ),
        ],
    )
    out = tmp_path / "out.pptx"
    render_deck(spec, TEMPLATE_PATH, _lookup(art), out)
    prs = Presentation(str(out))
    slide = prs.slides[0]
    table_shape = next((s for s in slide.shapes if s.has_table), None)
    assert table_shape is not None, "expected table shape"
    table = table_shape.table
    assert table.rows[0].cells[0].text == "Name"
    assert table.rows[0].cells[1].text == "Value"
    assert table.rows[1].cells[0].text == "row0"
    assert table.rows[3].cells[0].text == "row2"


def test_max_rows_truncation_adds_footer(tmp_path):
    art = _make_table_artifact(n_rows=20)
    spec = DeckSpec(
        title="Deck",
        slides=[
            TableSlide(
                type="table", title="Big table", artifact_id=art.id, max_rows=5
            ),
        ],
    )
    out = tmp_path / "out.pptx"
    render_deck(spec, TEMPLATE_PATH, _lookup(art), out)
    prs = Presentation(str(out))
    slide = prs.slides[0]
    table_shape = next((s for s in slide.shapes if s.has_table), None)
    assert table_shape is not None
    assert len(table_shape.table.rows) == 6  # 5 data + 1 header
    text_runs = [
        r.text
        for shape in slide.shapes
        if shape.has_text_frame
        for p in shape.text_frame.paragraphs
        for r in p.runs
    ]
    assert any("Showing 5 of 20 rows" in t for t in text_runs)


def test_validation_rejects_long_title():
    with pytest.raises(ValidationError):
        TitleSlide(type="title", title="x" * 200)


def test_validation_rejects_too_many_bullets():
    with pytest.raises(ValidationError):
        BulletsSlide(type="bullets", title="t", bullets=["a", "b", "c", "d", "e", "f"])


def test_validation_rejects_long_bullet():
    with pytest.raises(ValidationError):
        BulletsSlide(type="bullets", title="t", bullets=["x" * 200])


def test_render_rejects_missing_artifact(tmp_path):
    spec = DeckSpec(
        title="Deck",
        slides=[
            ChartSlide(type="chart", title="Missing", artifact_id="art_doesnotexist"),
        ],
    )
    with pytest.raises(PptxRenderError, match="not found"):
        render_deck(spec, TEMPLATE_PATH, _lookup(), tmp_path / "out.pptx")


def test_render_rejects_wrong_kind_artifact(tmp_path):
    text_art = _make_text_artifact()
    spec = DeckSpec(
        title="Deck",
        slides=[
            ChartSlide(type="chart", title="Wrong kind", artifact_id=text_art.id),
        ],
    )
    with pytest.raises(PptxRenderError, match="kind="):
        render_deck(spec, TEMPLATE_PATH, _lookup(text_art), tmp_path / "out.pptx")


def test_layout_lookup_missing(tmp_path):
    bad_template = tmp_path / "bad.pptx"
    from pptx import Presentation as _Pres

    p = _Pres()
    for layout in list(p.slide_layouts):
        layout.name = "junk"
    p.save(str(bad_template))
    spec = DeckSpec(
        title="Deck",
        slides=[TitleSlide(type="title", title="T")],
    )
    with pytest.raises(PptxRenderError, match="missing layout"):
        render_deck(spec, bad_template, _lookup(), tmp_path / "out.pptx")


def test_render_full_deck_smoke(tmp_path):
    img = _make_image_artifact("art_img_1")
    tbl = _make_table_artifact("art_tbl_1", n_rows=4)
    spec = DeckSpec(
        title="Q3 Review",
        slides=[
            TitleSlide(type="title", title="Q3 Review", subtitle="Top-line"),
            SectionSlide(type="section", title="Top-line", eyebrow="Q3 RESULTS"),
            ChartSlide(
                type="chart", title="Revenue", artifact_id=img.id, caption="+18% YoY"
            ),
            TableSlide(
                type="table", title="Top customers", artifact_id=tbl.id, max_rows=8
            ),
            BulletsSlide(
                type="bullets",
                title="Drivers",
                bullets=["Mid-market +27%", "Churn -2pp", "EU pipeline +40%"],
            ),
            ConclusionSlide(
                type="conclusion",
                title="Takeaways",
                bullets=["Scale pricing", "Hire EU SE", "Hold churn"],
                cta="Decision by 2026-05-15",
            ),
        ],
    )
    out = tmp_path / "out.pptx"
    render_deck(spec, TEMPLATE_PATH, _lookup(img, tbl), out)
    prs = Presentation(str(out))
    assert len(prs.slides) == 6
    assert out.stat().st_size > 5000
