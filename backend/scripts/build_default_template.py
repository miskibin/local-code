"""Build backend/data/pptx_templates/default.pptx with 6 named layouts.

Run once: `cd backend && uv run python scripts/build_default_template.py`.
Resulting file is committed; rerun only when the template needs to change.

Output template has 6 layouts named: title, section, bullets, chart, table,
conclusion. Layouts are intentionally empty (no placeholders) — the
renderer (`app.services.pptx_render`) draws all text and shapes per slide
at fixed coordinates. Layout names exist so the renderer can pick a layout
and so consumers can recognize a slide's purpose by inspecting its layout.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

OUT_PATH = Path(__file__).resolve().parents[1] / "data" / "pptx_templates" / "default.pptx"
LAYOUT_NAMES = ["title", "section", "bullets", "chart", "table", "conclusion"]
PAGE_W = Inches(13.333)
PAGE_H = Inches(7.5)


def _build():
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = PAGE_W
    prs.slide_height = PAGE_H

    master = prs.slide_master
    master.background.fill.solid()
    master.background.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")

    layouts = list(prs.slide_layouts)
    for layout, name in zip(layouts[: len(LAYOUT_NAMES)], LAYOUT_NAMES, strict=False):
        layout.name = name
        for ph in list(layout.placeholders):
            ph.element.getparent().remove(ph.element)

    prs.save(str(OUT_PATH))
    print(f"wrote {OUT_PATH} ({OUT_PATH.stat().st_size} bytes)")


if __name__ == "__main__":
    _build()
